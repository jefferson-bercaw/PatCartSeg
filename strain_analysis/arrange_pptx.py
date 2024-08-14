from pptx import Presentation
from pptx.util import Inches, Pt
import os


def add_vertical_texts(prs, slide, msgs):
    lefts = [0.1, 0.1, prs.slide_width/2, prs.slide_width/2]
    tops = [prs.slide_height * 0.25, prs.slide_height * 0.75, prs.slide_height * 0.25, prs.slide_height * 0.75]
    width = Inches(1)
    height = Inches(1)

    for i, msg in enumerate(msgs):
        if msg is not None:
            text_box = slide.shapes.add_textbox(
                lefts[i],
                tops[i],
                width,
                height
            )
            text_frame = text_box.text_frame
            text_frame.text = msg
            text_box.rotation = -90

    return prs, slide


if __name__ == "__main__":
    prs = Presentation()

    start_path = "R:\\DefratePrivate\\Bercaw\\Patella_Autoseg\\results\\cropped_7"

    subjs = os.listdir(start_path)
    subjs = [subj for subj in subjs if "." not in subj]

    for subj in subjs:
        img_names = []

        # Slide A: Averaging thickness maps
        img_names.append(os.path.join(start_path, subj, f"{subj}pre_thick_raw.png"))
        img_names.append(os.path.join(start_path, subj, f"{subj}pre_thick_avg.png"))
        img_names.append(os.path.join(start_path, subj, f"{subj}post_thick_raw.png"))
        img_names.append(os.path.join(start_path, subj, f"{subj}post_thick_avg.png"))

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        aspect_ratio = 1024 / 768
        slide_aspect_ratio = 2048 / 1536
        img_width = Inches(4)
        img_height = Inches(3)

        prs.slide_width = Inches(8.4)
        prs.slide_height = Inches(6.4)

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        positions = [
            (Inches(0.2), Inches(0.2)),  # Top-left corner
            (slide_width - img_width - Inches(0.2), Inches(0.2)),  # Top-right corner
            (Inches(0.2), slide_height - img_height - Inches(0.2)),  # Bottom-left corner
            (slide_width - img_width - Inches(0.2), slide_height - img_height - Inches(0.2))  # Bottom-right corner
        ]

        for i, img_path in enumerate(img_names):
            left, top = positions[i]
            slide.shapes.add_picture(img_path, left, top, width=img_width, height=img_height)

        # Add text to the top middle of the slide
        text_box_width = Inches(3)  # Adjust this to fit your text
        text_box_height = Inches(0.3)
        text_box_left = (slide_width - text_box_width) / 2
        text_box_top = slide_height - text_box_height

        text_box = slide.shapes.add_textbox(text_box_left, text_box_top, text_box_width, text_box_height)
        text_frame = text_box.text_frame
        text_frame.text = f"{subj} Thickness Averaging"  # Replace this with your desired text
        font = text_frame.paragraphs[0].font
        font.size = Pt(12)

        # Add vertical text
        prs, slide = add_vertical_texts(prs, slide, msgs=["Pre-thickness", "Post-thickness", "Pre-thickness_avg", "Post-thickness_avg"])

        ############ Slide B: Pre-post thicknesses to strain map and averaged strain map #############
        img_names = []

        img_names.append(os.path.join(start_path, subj, f"{subj}pre_thick_avg.png"))
        img_names.append(os.path.join(start_path, subj, f"{subj}_strain_raw.png"))
        img_names.append(os.path.join(start_path, subj, f"{subj}post_thick_avg.png"))
        img_names.append(os.path.join(start_path, subj, f"{subj}_strain_avg.png"))

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        for i, img_path in enumerate(img_names):
            left, top = positions[i]
            slide.shapes.add_picture(img_path, left, top, width=img_width, height=img_height)

        # Add text to the top middle of the slide
        text_box = slide.shapes.add_textbox(text_box_left, text_box_top, text_box_width, text_box_height)
        text_frame = text_box.text_frame
        text_frame.text = f"{subj} Strain Results"  # Replace this with your desired text
        font = text_frame.paragraphs[0].font
        font.size = Pt(12)

        prs, slide = add_vertical_texts(prs, slide, msgs=["Pre-thickness_avg", "Post-thickness_avg", "Raw Strain", "Averaged Strain"])

        # Slide C: Strain to cropped strain
        img_names = []

        img_names.append(os.path.join(start_path, subj, f"{subj}_strain_avg.png"))
        img_names.append(os.path.join(start_path, subj, f"{subj}_strain_boundaries_removed.png"))

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        for i, img_path in enumerate(img_names):
            left, top = positions[i+2]
            slide.shapes.add_picture(img_path, left, top, width=img_width, height=img_height)

        # Add text to the top middle of the slide
        text_box = slide.shapes.add_textbox(text_box_left, text_box_top, text_box_width, text_box_height)
        text_frame = text_box.text_frame
        text_frame.text = f"{subj} Strain Cropping"  # Replace this with your desired text
        font = text_frame.paragraphs[0].font
        font.size = Pt(12)

        prs, slide = add_vertical_texts(prs, slide, msgs=[None, "Averaged Strain", None, "Avg Strain Cropped"])

    prs.save("R:\\DefratePrivate\\Bercaw\\Patella_Autoseg\\results\\cropped_7\\results.pptx")
